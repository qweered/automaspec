import os
import re
import shutil
import tempfile
import zipfile
import xml.etree.ElementTree as ET

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt


NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}


def _int_attr(element: ET.Element, attr: str) -> int | None:
    value = element.get(attr)
    if value is None:
        return None
    return int(value)


def _get_transform(element: ET.Element) -> tuple[int, int, int, int] | None:
    xfrm = element.find('./p:spPr/a:xfrm', NS)
    if xfrm is None:
        xfrm = element.find('.//a:xfrm', NS)
    if xfrm is None:
        return None
    off = xfrm.find('./a:off', NS)
    ext = xfrm.find('./a:ext', NS)
    if off is None or ext is None:
        return None
    x = _int_attr(off, 'x')
    y = _int_attr(off, 'y')
    cx = _int_attr(ext, 'cx')
    cy = _int_attr(ext, 'cy')
    if x is None or y is None or cx is None or cy is None:
        return None
    return x, y, cx, cy


def _get_pic_transform(pic: ET.Element) -> tuple[int, int, int, int] | None:
    xfrm = pic.find('./p:spPr/a:xfrm', NS)
    if xfrm is None:
        xfrm = pic.find('.//a:xfrm', NS)
    if xfrm is None:
        return None
    off = xfrm.find('./a:off', NS)
    ext = xfrm.find('./a:ext', NS)
    if off is None or ext is None:
        return None
    x = _int_attr(off, 'x')
    y = _int_attr(off, 'y')
    cx = _int_attr(ext, 'cx')
    cy = _int_attr(ext, 'cy')
    if x is None or y is None or cx is None or cy is None:
        return None
    return x, y, cx, cy


def _shape_type_from_prst(prst: str | None) -> int | None:
    if prst is None:
        return None
    if prst in {'rect', 'roundRect'}:
        return MSO_SHAPE.ROUNDED_RECTANGLE if prst == 'roundRect' else MSO_SHAPE.RECTANGLE
    if prst in {'ellipse'}:
        return MSO_SHAPE.OVAL
    return None


def _rgb_from_srgb(element: ET.Element | None) -> RGBColor | None:
    if element is None:
        return None
    value = element.get('val')
    if value is None:
        return None
    return RGBColor.from_string(value)


def _apply_fill_and_line(shape, sp: ET.Element) -> None:
    sp_pr = sp.find('./p:spPr', NS)
    if sp_pr is None:
        return

    if sp_pr.find('./a:blipFill/a:blip', NS) is not None:
        shape.fill.background()
        ln = sp_pr.find('./a:ln', NS)
        if ln is not None:
            ln_fill = ln.find('./a:solidFill/a:srgbClr', NS)
            ln_no = ln.find('./a:noFill', NS)
            if ln_no is not None:
                shape.line.fill.background()
            elif ln_fill is not None:
                shape.line.color.rgb = _rgb_from_srgb(ln_fill)
            w = ln.get('w')
            if w is not None:
                shape.line.width = Emu(int(w))
        return

    fill = sp_pr.find('./a:solidFill/a:srgbClr', NS)
    no_fill = sp_pr.find('./a:noFill', NS)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb_from_srgb(fill)
    elif no_fill is not None:
        shape.fill.background()

    ln = sp_pr.find('./a:ln', NS)
    if ln is not None:
        ln_fill = ln.find('./a:solidFill/a:srgbClr', NS)
        ln_no = ln.find('./a:noFill', NS)
        if ln_no is not None:
            shape.line.fill.background()
        elif ln_fill is not None:
            shape.line.color.rgb = _rgb_from_srgb(ln_fill)
        w = ln.get('w')
        if w is not None:
            shape.line.width = Emu(int(w))


def _get_text_runs(sp: ET.Element) -> list[tuple[str, int | None, str | None, str | None]]:
    result: list[tuple[str, int | None, str | None, str | None]] = []
    tx_body = sp.find('./p:txBody', NS)
    if tx_body is None:
        return result

    for para in tx_body.findall('./a:p', NS):
        para_text_parts: list[str] = []
        font_size: int | None = None
        align: str | None = None
        bullet: str | None = None

        ppr = para.find('./a:pPr', NS)
        if ppr is not None:
            align = ppr.get('algn')
            if ppr.find('./a:buChar', NS) is not None or ppr.find('./a:buAutoNum', NS) is not None:
                bullet = '•'

        for run in para.findall('./a:r', NS):
            rpr = run.find('./a:rPr', NS)
            if rpr is not None and font_size is None:
                sz = rpr.get('sz')
                if sz is not None:
                    font_size = int(sz)
            for t in run.findall('./a:t', NS):
                if t.text:
                    para_text_parts.append(t.text)

        for t in para.findall('./a:t', NS):
            if t.text:
                para_text_parts.append(t.text)

        text = ''.join(para_text_parts).strip()
        if text:
            result.append((text, font_size, align, bullet))

    return result


def _apply_text(shape, sp: ET.Element) -> None:
    runs = _get_text_runs(sp)
    if not runs:
        return

    tf = shape.text_frame
    tf.clear()

    for i, (text, font_size, align, bullet) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'{bullet} {text}' if bullet else text

        if align == 'ctr':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'r':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        if font_size is not None:
            p.font.size = Pt(font_size / 100)


def rebuild_editable_pptx(src_pptx: str, out_pptx: str) -> None:
    temp_dir = tempfile.mkdtemp(prefix='pptx_rebuild_')
    extract_dir = os.path.join(temp_dir, 'pptx')
    media_dir = os.path.join(temp_dir, 'media')
    os.makedirs(extract_dir, exist_ok=True)
    os.makedirs(media_dir, exist_ok=True)

    try:
        with zipfile.ZipFile(src_pptx, 'r') as z:
            z.extractall(extract_dir)
            for name in z.namelist():
                if name.startswith('ppt/media/'):
                    target = os.path.join(media_dir, os.path.basename(name))
                    with z.open(name) as fsrc, open(target, 'wb') as fdst:
                        shutil.copyfileobj(fsrc, fdst)

        prs = Presentation()
        prs.slide_width = Emu(960 * 9525)
        prs.slide_height = Emu(540 * 9525)

        pres_xml = os.path.join(extract_dir, 'ppt', 'presentation.xml')
        if os.path.exists(pres_xml):
            root = ET.parse(pres_xml).getroot()
            sld_sz = root.find('./p:sldSz', NS)
            if sld_sz is not None:
                cx = _int_attr(sld_sz, 'cx')
                cy = _int_attr(sld_sz, 'cy')
                if cx and cy:
                    prs.slide_width = Emu(cx)
                    prs.slide_height = Emu(cy)

        slide_paths = []
        slides_dir = os.path.join(extract_dir, 'ppt', 'slides')
        for name in os.listdir(slides_dir):
            m = re.match(r'slide(\d+)\.xml$', name)
            if m:
                slide_paths.append((int(m.group(1)), os.path.join(slides_dir, name)))
        slide_paths.sort()

        for idx, slide_path in slide_paths:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            rels_map: dict[str, str] = {}
            rels_path = os.path.join(extract_dir, 'ppt', 'slides', '_rels', f'slide{idx}.xml.rels')
            if os.path.exists(rels_path):
                rels_root = ET.parse(rels_path).getroot()
                for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship'):
                    rid = rel.get('Id')
                    target = rel.get('Target')
                    if rid and target and target.startswith('../media/'):
                        rels_map[rid] = os.path.join(media_dir, os.path.basename(target))

            root = ET.parse(slide_path).getroot()

            bg_blip = root.find('.//p:bg/p:bgPr/a:blipFill/a:blip', NS)
            if bg_blip is not None:
                rid = bg_blip.get('{%s}embed' % NS['r'])
                if rid is not None:
                    img_path = rels_map.get(rid)
                    if img_path and os.path.exists(img_path):
                        slide.shapes.add_picture(
                            img_path,
                            Emu(0),
                            Emu(0),
                            prs.slide_width,
                            prs.slide_height,
                        )

            for pic in root.findall('.//p:pic', NS):
                blip = pic.find('.//a:blip', NS)
                if blip is None:
                    continue
                rid = blip.get('{%s}embed' % NS['r'])
                if rid is None:
                    continue
                img_path = rels_map.get(rid)
                if not img_path or not os.path.exists(img_path):
                    continue
                transform = _get_pic_transform(pic)
                if transform is None:
                    continue
                x, y, cx, cy = transform
                slide.shapes.add_picture(img_path, Emu(x), Emu(y), Emu(cx), Emu(cy))

            for sp in root.findall('.//p:sp', NS):
                transform = _get_transform(sp)
                if transform is None:
                    continue
                x, y, cx, cy = transform

                blip = sp.find('./p:spPr/a:blipFill/a:blip', NS)
                if blip is not None:
                    rid = blip.get('{%s}embed' % NS['r'])
                    if rid is not None:
                        img_path = rels_map.get(rid)
                        if img_path and os.path.exists(img_path):
                            slide.shapes.add_picture(img_path, Emu(x), Emu(y), Emu(cx), Emu(cy))

                prst = None
                geom = sp.find('./p:spPr/a:prstGeom', NS)
                if geom is not None:
                    prst = geom.get('prst')

                shape_type = _shape_type_from_prst(prst)
                if shape_type is None:
                    shape = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
                else:
                    shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(cx), Emu(cy))
                    _apply_fill_and_line(shape, sp)

                runs = _get_text_runs(sp)
                if runs:
                    raw_text = ' '.join([t for (t, _, _, _) in runs]).strip()
                    if not re.search(r'Image failed to load:', raw_text, re.IGNORECASE):
                        _apply_text(shape, sp)

        out_dir = os.path.dirname(out_pptx)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)
        prs.save(out_pptx)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def rebuild_pdf_overlay_pptx(src_pptx: str, src_pdf: str, out_pptx: str) -> None:
    temp_dir = tempfile.mkdtemp(prefix='pptx_overlay_')
    extract_dir = os.path.join(temp_dir, 'pptx')
    os.makedirs(extract_dir, exist_ok=True)

    render_dir = os.path.join(temp_dir, 'render')
    os.makedirs(render_dir, exist_ok=True)

    try:
        with zipfile.ZipFile(src_pptx, 'r') as z:
            z.extractall(extract_dir)

        prs = Presentation()
        prs.slide_width = Emu(960 * 9525)
        prs.slide_height = Emu(540 * 9525)

        pres_xml = os.path.join(extract_dir, 'ppt', 'presentation.xml')
        if os.path.exists(pres_xml):
            root = ET.parse(pres_xml).getroot()
            sld_sz = root.find('./p:sldSz', NS)
            if sld_sz is not None:
                cx = _int_attr(sld_sz, 'cx')
                cy = _int_attr(sld_sz, 'cy')
                if cx and cy:
                    prs.slide_width = Emu(cx)
                    prs.slide_height = Emu(cy)

        pdf = fitz.open(src_pdf)
        slide_paths: list[tuple[int, str]] = []
        slides_dir = os.path.join(extract_dir, 'ppt', 'slides')
        for name in os.listdir(slides_dir):
            m = re.match(r'slide(\d+)\.xml$', name)
            if m:
                slide_paths.append((int(m.group(1)), os.path.join(slides_dir, name)))
        slide_paths.sort()

        for idx, slide_path in slide_paths:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            page_index = idx - 1
            if 0 <= page_index < len(pdf):
                pix = pdf[page_index].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                bg_path = os.path.join(render_dir, f'page-{idx:02d}.png')
                pix.save(bg_path)
                slide.shapes.add_picture(bg_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

            root = ET.parse(slide_path).getroot()
            for sp in root.findall('.//p:sp', NS):
                transform = _get_transform(sp)
                if transform is None:
                    continue
                x, y, cx, cy = transform
                runs = _get_text_runs(sp)
                if not runs:
                    continue
                raw_text = ' '.join([t for (t, _, _, _) in runs]).strip()
                if re.search(r'Image failed to load:', raw_text, re.IGNORECASE):
                    continue
                shape = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(cx), Emu(cy))
                shape.fill.background()
                shape.line.fill.background()
                _apply_text(shape, sp)

        out_dir = os.path.dirname(out_pptx)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)
        prs.save(out_pptx)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


if __name__ == '__main__':
    repo_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
    src = os.path.join(repo_root, 'test-presentation', 'Automaspec Presentation.pptx')
    pdf = os.path.join(repo_root, 'test-presentation', 'Automaspec Presentation.pdf')
    out = os.path.join(repo_root, 'test-presentation', 'Automaspec Presentation.fixed-editable.pptx')
    rebuild_editable_pptx(src, out)
    print(out)
